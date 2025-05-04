import streamlit as st
import tempfile
import whisperx
from fpdf import FPDF
import google.generativeai as genai
from pptx import Presentation
from datetime import datetime

# --- Page Configuration ---
st.set_page_config(
    page_title="Voice AI Assistant",
    page_icon="üé§",
    layout="centered",
    initial_sidebar_state="collapsed"
)

# --- Gemini Configuration ---
from dotenv import load_dotenv
import os

load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))


model = genai.GenerativeModel('gemini-1.5-flash')
# --- Session State Management ---
session_defaults = {
    "page": "main",
    "transcription": "",
    "summary": "",
    "audio_data": None,
    "audio_path": "",
    "messages": [],
    "ppt_path": None,
    "summary_path": None,
    "ppt_title": "",
    "ppt_headings": ""
}

for key, value in session_defaults.items():
    if key not in st.session_state:
        st.session_state[key] = value

# --- Helper Functions ---
def transcribe_audio():
    with st.spinner("Transcribing audio..."):
        model = whisperx.load_model("base.en", device="cpu", compute_type="float32")
        result = model.transcribe(st.session_state.audio_path)
        alignment_model, metadata = whisperx.load_align_model("en", device="cpu")
        audio = whisperx.load_audio(st.session_state.audio_path)
        aligned_result = whisperx.align(
            result["segments"],
            alignment_model,
            metadata,
            audio,
            device="cpu"
        )
        return "\n".join(
            f"[{seg['start']:.2f}s] {seg['text']}" 
            for seg in aligned_result["segments"]
        )

def generate_summary():
    try:
        prompt = f"""
        Create a professional summary from this audio transcript:
        {st.session_state.transcription}
        
        Include:
        - Key discussion points
        - Important figures/dates
        - Action items
        - Recommendations
        - Overall sentiment analysis
        
        Format with markdown headings and bullet points.
        """
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"‚ùå Error generating summary: {str(e)}"

def create_presentation(title, headings):
    try:
        prs = Presentation()
        
        # Title Slide
        title_slide = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = f"Generated {datetime.now().strftime('%d %b %Y %H:%M')}"
        
        # Process Headings
        for heading in [h.strip() for h in headings.split('\n') if h.strip()]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = heading
            
            # Generate content
            prompt = f"""
            Create 3-5 bullet points for '{heading}' using:
            {st.session_state.transcription}
            
            Rules:
            - Use concise business language
            - Include key figures/dates
            - Format as markdown list
            """
            response = model.generate_content(prompt)
            content.text = '\n'.join([f"- {point}" for point in response.text.split('\n') if point.strip()])
        
        # Save PPT
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            prs.save(tmp_file.name)
            return tmp_file.name
            
    except Exception as e:
        st.error(f"PPT Error: {str(e)}")
        return None

def chat_response(prompt):
    system_prompt = f"""
    ROLE: Voice Analysis Assistant
    CONTEXT: {st.session_state.transcription}
    HISTORY: {st.session_state.messages[-5:]}
    
    Rules:
    1. Answer strictly based on the audio content
    2. If question is unrelated to audio, respond politely with different  wordings for each unrelated questions
    3. Maintain conversational flow
    4.You are a chatbot named "Chat with your voice" . You need to understand the audio trabscribe intent or meaning and the query meaning and respond correctly
    5. Use markdown formatting when appropriate
    6. Keep responses concise but helpful
    7. Use markdown formatting (**bold**, *italics*, lists) when helpful
    8. Never mention you're an AI or language model
    9. Maintain friendly, professional tone but dont need to wish or greeting in each response
    10.Acknowledge previous interactions when relevant
    11.In some cases for unrelated questions you can also give like  "I specialize in analyzing the current recording. Ask me about: [list key topics]"
    
    USER QUESTION: {prompt}
    """
    
    try:
        response = model.generate_content(system_prompt)
        return response.text
    except Exception as e:
        return f"‚ö†Ô∏è Error: {str(e)}"

# --- Page Handlers ---
def main_page():
    st.title("üé§ Voice AI Assistant")
    
    # Audio Recording Section
    col1, col2 = st.columns([3, 1])
    with col1:
        if st.session_state.audio_data is None:
            audio_input = st.audio_input("Start Recording", key="recorder")
        else:
            st.audio(st.session_state.audio_data, format="audio/wav")
    
    with col2:
        if st.session_state.audio_data and st.button("üîÑ New Session"):
            st.session_state.audio_data = None
            st.session_state.audio_path = ""
            st.session_state.transcription = ""
            st.session_state.messages = []
            st.rerun()
    
    # Process Recording
    if st.session_state.audio_data is None and 'recorder' in st.session_state:
        if st.session_state.recorder is not None:
            st.session_state.audio_data = st.session_state.recorder
            with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_audio:
                tmp_audio.write(st.session_state.audio_data.read())
                st.session_state.audio_path = tmp_audio.name
    
    # Transcription
    if st.session_state.audio_data and not st.session_state.transcription:
        st.session_state.transcription = transcribe_audio()
    
    if st.session_state.transcription:
        with st.expander("üìú View Transcript"):
            st.write(st.session_state.transcription)
        
        # Action Buttons
        st.markdown("---")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            if st.button("üìÑ Generate Summary"):
                st.session_state.page = "summary"
                st.rerun()
        
        with col2:
            if st.button("üìä Create PPT"):
                st.session_state.page = "ppt"
                st.rerun()
        
        with col3:
            if st.button("üí¨ Chat with Audio"):
                st.session_state.page = "chat"
                st.rerun()

def summary_page():
    st.title("üìÑ Executive Summary")
    
    if st.button("‚Üê Back to Main"):
        st.session_state.page = "main"
        st.rerun()
    
    if not st.session_state.summary:
        with st.spinner("Generating professional summary..."):
            st.session_state.summary = generate_summary()
            
            # Save summary
            with tempfile.NamedTemporaryFile(delete=False, suffix=".txt") as tmp_file:
                tmp_file.write(st.session_state.summary.encode())
                st.session_state.summary_path = tmp_file.name
    
    st.markdown(st.session_state.summary)
    
    st.download_button(
        "‚¨áÔ∏è Download Summary",
        open(st.session_state.summary_path, "rb"),
        "voice_summary.md",
        "text/markdown"
    )

def ppt_page():
    st.title("üìä PPT Generator")
    
    if st.button("‚Üê Back to Main"):
        st.session_state.page = "main"
        st.rerun()
    
    with st.form("ppt_config"):
        st.subheader("Presentation Settings")
        title = st.text_input("Presentation Title", "Voice Analysis Report")
        headings = st.text_area("Slide Headings (one per line)", 
                              "Introduction\nKey Findings\nData Analysis\nRecommendations\nNext Steps")
        
        if st.form_submit_button("üöÄ Generate Presentation"):
            with st.spinner("Building professional slides..."):
                ppt_path = create_presentation(title, headings)
                st.session_state.ppt_path = ppt_path
    
    if st.session_state.ppt_path:
        st.success("‚úÖ Presentation Ready!")
        st.download_button(
            "‚¨áÔ∏è Download PowerPoint",
            open(st.session_state.ppt_path, "rb"),
            "voice_analysis.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

def chat_page():
    st.title("üí¨ Voice Analysis Chat")
    
    if st.button("‚Üê Back to Main"):
        st.session_state.page = "main"
        st.rerun()
    
    # Chat History
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    
    # Chat Input
    if prompt := st.chat_input("Ask about the recording..."):
        # Add user message
        st.session_state.messages.append({"role": "user", "content": prompt})
        
        # Generate response
        with st.spinner("Analyzing..."):
            response = chat_response(prompt)
            st.session_state.messages.append({"role": "assistant", "content": response})
            st.rerun()

# --- App Flow ---
if st.session_state.page == "main":
    main_page()
elif st.session_state.page == "summary":
    summary_page()
elif st.session_state.page == "ppt":
    ppt_page()
elif st.session_state.page == "chat":
    chat_page()
else:
    st.session_state.page = "main"
    st.rerun()



